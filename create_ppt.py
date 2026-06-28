from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
PRIMARY = RGBColor(0xD8, 0x40, 0x40)
SECONDARY = RGBColor(0x8E, 0x16, 0x16)
ACCENT = RGBColor(0x1D, 0x16, 0x16)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
MID_GRAY = RGBColor(0x9C, 0xA3, 0xAF)
DARK_GRAY = RGBColor(0x4B, 0x55, 0x63)
TABLE_HEADER_BG = RGBColor(0xD8, 0x40, 0x40)
TABLE_ALT_BG = RGBColor(0xFE, 0xF2, 0xF2)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color, shape_type=MSO_SHAPE.RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18, color=ACCENT, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_textbox(slide, left, top, width, height, items, font_size=16, color=ACCENT, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = spacing
        p.level = 0
    return txBox


def add_table(slide, left, top, width, height, rows, cols, data, col_widths=None):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c]
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.name = 'Calibri'
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                    paragraph.alignment = PP_ALIGN.CENTER
                else:
                    paragraph.font.color.rgb = ACCENT
            
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER_BG
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT_BG
    
    return table_shape


def add_header_bar(slide):
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), PRIMARY)


def add_footer(slide, text="LiquidPedia | ABP - Telkom University"):
    add_shape(slide, Inches(0), Inches(7.0), SLIDE_W, Inches(0.5), ACCENT)
    add_textbox(slide, Inches(0.5), Inches(7.05), Inches(12), Inches(0.4), text, 
                font_size=10, color=RGBColor(0x9C, 0xA3, 0xAF), alignment=PP_ALIGN.CENTER)


def add_section_number(slide, number, text):
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(2), Inches(0.5), number,
                font_size=14, color=PRIMARY, bold=True)
    add_textbox(slide, Inches(0.5), Inches(0.7), Inches(12), Inches(0.6), text,
                font_size=28, color=ACCENT, bold=True)


# ===================== SLIDE 1: JUDUL =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)
add_shape(slide, Inches(0), Inches(0), Inches(5.5), SLIDE_H, PRIMARY)

add_textbox(slide, Inches(6.5), Inches(1.5), Inches(6), Inches(1.2),
            "LiquidPedia", font_size=54, color=PRIMARY, bold=True)
add_textbox(slide, Inches(6.5), Inches(2.7), Inches(6), Inches(0.8),
            "Web E-Commerce Liquid & Vape", font_size=28, color=ACCENT)
add_textbox(slide, Inches(6.5), Inches(4.0), Inches(6), Inches(0.5),
            "Analisis dan Perancangan Perangkat Lunak", font_size=16, color=MID_GRAY)
add_textbox(slide, Inches(6.5), Inches(4.5), Inches(6), Inches(0.5),
            "Informatika — Telkom University", font_size=16, color=MID_GRAY)

add_shape(slide, Inches(6.5), Inches(5.5), Inches(2), Inches(0.06), PRIMARY)
add_textbox(slide, Inches(6.5), Inches(5.8), Inches(6), Inches(0.5),
            "Dosen: ...", font_size=14, color=DARK_GRAY)
add_textbox(slide, Inches(6.5), Inches(6.2), Inches(6), Inches(0.5),
            "Kelompok: ...", font_size=14, color=DARK_GRAY)


# ===================== SLIDE 2: ALASAN =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)
add_footer(slide)
add_section_number(slide, "01", "Alasan / Latar Belakang")

data = [
    ["Masalah", "Solusi"],
    ["Penjualan liquid & vape masih manual\n(WhatsApp / offline)", "Platform e-commerce web khusus\nliquid & vape"],
    ["Pelanggan sulit membandingkan\nproduk dan harga", "Katalog produk dengan filter\nkategori + search bar"],
    ["Tidak ada sistem pemesanan\nyang terstruktur", "Keranjang belanja session-based\n+ checkout dengan location picker"],
    ["Admin kesulitan mengelola\nproduk dan pesanan", "Panel admin dengan CRUD produk\n+ dashboard statistik"],
]
add_table(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(4), 5, 2, data,
          col_widths=[Inches(5.5), Inches(6.2)])

add_textbox(slide, Inches(0.8), Inches(5.8), Inches(11.7), Inches(1),
            "Tujuan: Menyediakan platform e-commerce khusus liquid dan vape berbasis web dengan dua peran pengguna (Customer & Admin)",
            font_size=14, color=DARK_GRAY)


# ===================== SLIDE 3: TECH STACK =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)
add_footer(slide)
add_section_number(slide, "02", "Tech Stack")

data = [
    ["Layer", "Teknologi"],
    ["Backend", "Laravel 13, PHP 8.3"],
    ["Database", "MySQL / MariaDB"],
    ["Frontend", "Tailwind CSS v4, Vite, Blade Templating"],
    ["Maps", "Leaflet.js + OpenStreetMap + Nominatim (gratis, tanpa API key)"],
    ["Auth", "Laravel Auth (session-based)"],
    ["Cart", "Session Database"],
    ["Arsitektur", "MVC (Model-View-Controller)"],
]
add_table(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(4.5), 8, 2, data,
          col_widths=[Inches(3.5), Inches(8.2)])


# ===================== SLIDE 4: AKTOR =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)
add_footer(slide)
add_section_number(slide, "03", "Aktor")

data = [
    ["Aktor", "Peran"],
    ["Customer\n(Pelanggan)", "Registrasi, login, melihat & menyaring produk,\nmengelola keranjang, checkout dengan location picker,\nkonfirmasi pembayaran, riwayat pesanan"],
    ["Admin\n(Pengelola Toko)", "Login panel khusus (is_admin = true),\ndashboard statistik, CRUD produk & kategori"],
]
add_table(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(3), 3, 2, data,
          col_widths=[Inches(3.5), Inches(8.2)])


# ===================== SLIDE 5: USE CASE DIAGRAM =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)
add_footer(slide)
add_section_number(slide, "04", "Use Case Diagram")

data = [
    ["Kode", "Use Case", "Aktor"],
    ["UC-01", "Registrasi Akun", "Customer"],
    ["UC-02", "Login Akun", "Customer"],
    ["UC-03", "Melihat & Menyaring Produk", "Customer"],
    ["UC-04", "Melihat Detail Produk", "Customer"],
    ["UC-05", "Mengelola Keranjang Belanja", "Customer"],
    ["UC-06", "Melakukan Checkout", "Customer"],
    ["UC-07", "Melihat Konfirmasi Pembayaran", "Customer"],
    ["UC-08", "Melihat Riwayat Pesanan", "Customer"],
    ["UC-09", "Login & Dashboard Admin", "Admin"],
    ["UC-10", "Mengelola Produk", "Admin"],
]
add_table(slide, Inches(0.8), Inches(1.5), Inches(7), Inches(5), 11, 3, data,
          col_widths=[Inches(1.2), Inches(3.8), Inches(2.0)])

add_textbox(slide, Inches(8.2), Inches(1.5), Inches(4.5), Inches(1.5),
            "Relasi <<include>>:", font_size=18, color=PRIMARY, bold=True)

relasi_items = [
    "UC-06 <<include>> UC-02",
    "UC-07 <<include>> UC-06",
    "UC-08 <<include>> UC-02",
    "UC-10 <<include>> UC-09",
]
add_bullet_textbox(slide, Inches(8.2), Inches(3.0), Inches(4.5), Inches(3.5), relasi_items, font_size=16)


# ===================== SLIDE 6: FUNCTIONAL REQUIREMENTS =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)
add_footer(slide)
add_section_number(slide, "05", "Functional Requirements")

data = [
    ["Kode", "Nama", "Aktor", "Deskripsi"],
    ["FR-01", "Registrasi Akun", "Customer", "Form registrasi, validasi input, hash password"],
    ["FR-02", "Login Akun", "Customer", "Autentikasi session-based"],
    ["FR-03", "Melihat & Menyaring Produk", "Customer", "Grid produk + filter kategori + search bar"],
    ["FR-04", "Melihat Detail Produk", "Customer", "Info lengkap: gambar, harga, deskripsi, badge"],
    ["FR-05", "Mengelola Keranjang", "Customer", "Tambah/ubah/hapus via session"],
    ["FR-06", "Melakukan Checkout", "Customer", "Form + location picker + pilih metode bayar"],
    ["FR-07", "Konfirmasi Pembayaran", "Customer", "Instruksi bayar sesuai metode"],
    ["FR-08", "Riwayat Pesanan", "Customer", "Daftar pesanan terurut"],
    ["FR-09", "Login & Dashboard Admin", "Admin", "Login khusus + statistik toko"],
    ["FR-10", "Mengelola Produk", "Admin", "CRUD produk + upload gambar"],
]
add_table(slide, Inches(0.3), Inches(1.5), Inches(12.7), Inches(5.5), 11, 4, data,
          col_widths=[Inches(0.9), Inches(2.5), Inches(1.3), Inches(8.0)])


# ===================== SLIDE 7: CLASS DIAGRAM =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)
add_footer(slide)
add_section_number(slide, "06", "Class Diagram")

data = [
    ["Entity", "Key Attributes", "Relasi"],
    ["User", "id, name, email, password, is_admin", "1:N → Order"],
    ["Category", "id, name, slug", "1:N → Product"],
    ["Product", "id, name, description, price,\ncategory_id, image, is_best_seller,\nis_new_arrival", "N:1 ← Category\n1:N → OrderItem"],
    ["Order", "id, user_id, order_number,\nshipping_*, payment_method,\npayment_status, total", "N:1 ← User\n1:N → OrderItem"],
    ["OrderItem", "id, order_id, product_id,\nproduct_name, quantity, price,\nsubtotal", "N:1 ← Order\nN:1 ← Product"],
]
add_table(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(4.5), 6, 3, data,
          col_widths=[Inches(2.0), Inches(5.5), Inches(4.8)])

add_textbox(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.5),
            "Cart: Session-based (disimpan di tabel sessions Laravel, tanpa entity terpisah)",
            font_size=13, color=DARK_GRAY, bold=True)


# ===================== SLIDE 8: SEQUENCE DIAGRAM =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)
add_footer(slide)
add_section_number(slide, "07", "Sequence Diagram Terpenting")

add_textbox(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.5),
            "1. UC-06: Melakukan Checkout (paling kompleks)", font_size=20, color=PRIMARY, bold=True)

checkout_items = [
    'Actor -> Form Checkout: Isi data pengiriman + pilih lokasi di map (Leaflet)',
    'Form -> Controller POST /checkout: processOrder(request)',
    'Controller -> Model: validasi input',
    'Controller -> Model: User::find(auth()->id)',
    'Controller -> Model: Order::create($data)',
    'Controller -> Model: loop tiap item session cart:',
    '    Model: OrderItem::create($item)',
    'Model -> DB: INSERT orders + order_items',
    'Controller -> Session: hapus cart',
    'Controller -> View: redirect ke payment-confirmation',
]
add_bullet_textbox(slide, Inches(0.5), Inches(1.9), Inches(5.8), Inches(4.5), checkout_items, font_size=13, spacing=Pt(4))

add_textbox(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(0.5),
            "2. UC-10: Mengelola Produk (Admin)", font_size=20, color=PRIMARY, bold=True)

admin_items = [
    'Actor -> Form: Isi data + upload gambar',
    'Form -> Controller POST /admin/products: store(request)',
    'Controller -> Model: validasi request',
    'Controller -> Storage: upload gambar ke storage',
    'Model -> DB: INSERT products',
    'Controller -> View: redirect ke daftar produk',
]
add_bullet_textbox(slide, Inches(6.8), Inches(1.9), Inches(6), Inches(2.5), admin_items, font_size=13, spacing=Pt(4))

add_textbox(slide, Inches(6.8), Inches(4.5), Inches(6), Inches(0.5),
            "3. UC-03: Melihat & Menyaring Produk", font_size=20, color=PRIMARY, bold=True)

filter_items = [
    'Actor -> Browser: Buka /products',
    'Browser -> Controller GET /products: index(request)',
    'Controller -> Model: Product::with("category")',
    'Model -> DB: SELECT + WHERE category = ?',
    '           + WHERE name LIKE %search%',
    'Controller -> View: products.blade.php',
]
add_bullet_textbox(slide, Inches(6.8), Inches(5.0), Inches(6), Inches(2.0), filter_items, font_size=13, spacing=Pt(3))


# ===================== SLIDE 9: ARSITEKTUR MVC =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)
add_footer(slide)
add_section_number(slide, "08", "Arsitektur MVC")

# Draw boxes for MVC flow
boxes = [
    (Inches(0.5), Inches(2.0), "Browser\n(Blade + Tailwind)", PRIMARY),
    (Inches(3.0), Inches(2.0), "Router\nweb.php", SECONDARY),
    (Inches(5.5), Inches(2.0), "Controller\n(Logic)", RGBColor(0x25, 0x63, 0xEB)),
    (Inches(8.0), Inches(2.0), "Model\n(Eloquent ORM)", RGBColor(0x05, 0x9C, 0x69)),
    (Inches(10.5), Inches(2.0), "Database\n(MySQL)", RGBColor(0x7C, 0x3A, 0xED)),
]

for left, top, text, color in boxes:
    shape = add_shape(slide, left, top, Inches(2.2), Inches(1.2), color, MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(10)

# Arrows between boxes (simple text arrows)
arrow_y = Inches(2.5)
add_textbox(slide, Inches(2.7), arrow_y, Inches(0.5), Inches(0.5), "→", font_size=24, color=DARK_GRAY, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(5.2), arrow_y, Inches(0.5), Inches(0.5), "→", font_size=24, color=DARK_GRAY, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(7.7), arrow_y, Inches(0.5), Inches(0.5), "→", font_size=24, color=DARK_GRAY, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(10.2), arrow_y, Inches(0.5), Inches(0.5), "→", font_size=24, color=DARK_GRAY, bold=True, alignment=PP_ALIGN.CENTER)

# Return arrow
add_textbox(slide, Inches(10.5), Inches(3.5), Inches(2.2), Inches(0.5), "↑\nResponse", font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2), "Alur:\n1. Browser mengirim HTTP request ke Laravel\n2. Router (web.php) mencocokkan URL ke Controller\n3. Controller memproses logika bisnis\n4. Model berinteraksi dengan database via Eloquent ORM\n5. Controller merender View (Blade + Tailwind) → dikirim ke Browser",
            font_size=14, color=DARK_GRAY)


# ===================== SLIDE 10: DATABASE SCHEMA =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)
add_footer(slide)
add_section_number(slide, "09", "Database Schema")

data = [
    ["Tabel", "Fungsi", "Relasi"],
    ["users", "Data customer & admin\n(flag is_admin)", "PK di order.user_id"],
    ["categories", "Kategori produk\n(Vape / Liquid)", "PK di product.category_id"],
    ["products", "Produk: gambar, harga, deskripsi,\nbadge best_seller / new_arrival", "FK → categories\nPK di order_item.product_id"],
    ["orders", "Pesanan + data pengiriman\n+ koordinat latitude/longitude", "FK → users\nPK di order_item.order_id"],
    ["order_items", "Item per pesanan (snapshot\nnama & harga saat checkout)", "FK → orders\nFK → products"],
    ["sessions", "Cart session\n(Laravel session database)", "—"],
]
add_table(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5), 7, 3, data,
          col_widths=[Inches(2.0), Inches(5.0), Inches(5.3)])


# ===================== SLIDE 11: UI SCREENS =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)
add_footer(slide)
add_section_number(slide, "10", "Halaman UI")

add_textbox(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(0.5),
            "Customer (8 halaman)", font_size=20, color=PRIMARY, bold=True)

customer_pages = [
    "1. Beranda (/) — Hero, Best Seller, New Arrival",
    "2. Katalog (/products) — Grid + filter + search",
    "3. Detail Produk (/products/{id})",
    "4. Keranjang (/cart) — Daftar item + total",
    "5. Checkout (/checkout) — Form + Leaflet map",
    "6. Konfirmasi Bayar (/orders/{id}/payment-confirmation)",
    "7. Riwayat Pesanan (/orders)",
    "8. Login / Register",
]
add_bullet_textbox(slide, Inches(0.5), Inches(1.9), Inches(6), Inches(4.5), customer_pages, font_size=14, spacing=Pt(4))

add_textbox(slide, Inches(7), Inches(1.3), Inches(6), Inches(0.5),
            "Admin (4 halaman)", font_size=20, color=PRIMARY, bold=True)

admin_pages = [
    "1. Login Admin (/admin/login)",
    "2. Dashboard (/admin) — Statistik toko",
    "3. Daftar Produk (/admin/products)",
    "4. Tambah / Edit Produk",
]
add_bullet_textbox(slide, Inches(7), Inches(1.9), Inches(6), Inches(3), admin_pages, font_size=14, spacing=Pt(4))

add_shape(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.06), PRIMARY)
add_textbox(slide, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.5),
            "Fitur Unggulan: Location Picker Leaflet.js + OpenStreetMap pada halaman Checkout",
            font_size=14, color=PRIMARY, bold=True)


# ===================== SLIDE 12: DEMO =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)
add_footer(slide)
add_section_number(slide, "11", "Demo")

data = [
    ["No", "Skenario", "Durasi"],
    ["1", "Home — Best seller & new arrival", "30 detik"],
    ["2", "Katalog + Search — Filter Vape/Liquid, cari produk", "45 detik"],
    ["3", "Detail Produk — Klik produk, lihat info lengkap", "30 detik"],
    ["4", "Register + Login — Buat akun baru", "45 detik"],
    ["5", "Keranjang — Tambah 2 produk, ubah quantity", "45 detik"],
    ["6", "Checkout — Location picker Leaflet ★", "1 menit"],
    ["7", "Konfirmasi Bayar — Instruksi transfer / QRIS", "30 detik"],
    ["8", "Admin Panel — Login, dashboard, tambah produk", "1 menit"],
]
add_table(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(4.5), 9, 3, data,
          col_widths=[Inches(0.8), Inches(8.5), Inches(2.4)])

add_shape(slide, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.06), PRIMARY)
add_textbox(slide, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.5),
            "★ Highlight: Location picker Leaflet.js — fitur paling unik dari aplikasi",
            font_size=16, color=PRIMARY, bold=True)


# ===================== SLIDE 13: PENUTUP =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)
add_footer(slide)

add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
            "Kesimpulan", font_size=32, color=PRIMARY, bold=True)
add_shape(slide, Inches(0.5), Inches(1.0), Inches(3), Inches(0.06), PRIMARY)

conclusion_items = [
    "Aplikasi e-commerce liquid & vape berbasis Laravel web",
    "10 use cases mencakup kebutuhan Customer dan Admin",
    "MVC architecture dengan session-based cart",
    "Fitur unggulan: Leaflet location picker, filter + search, CRUD produk",
]
add_bullet_textbox(slide, Inches(0.5), Inches(1.5), Inches(12), Inches(3), conclusion_items, font_size=18, spacing=Pt(10))

add_shape(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.06), LIGHT_GRAY)

add_textbox(slide, Inches(0.5), Inches(4.8), Inches(12), Inches(0.5),
            "Link Repository", font_size=22, color=ACCENT, bold=True)
add_textbox(slide, Inches(0.5), Inches(5.3), Inches(12), Inches(0.5),
            "github.com/notsorcerer/ABP-WEB", font_size=18, color=PRIMARY)

add_textbox(slide, Inches(0.5), Inches(6.0), Inches(12), Inches(0.5),
            "Terima Kasih  |  Q & A", font_size=28, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)


# Save
output_path = "C:\\Kuliah\\Semester6\\ABP\\Tubes\\liquid\\Presentasi_LiquidPedia.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
